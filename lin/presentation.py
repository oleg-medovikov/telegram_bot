import os,datetime,sqlalchemy
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Inches
import pandas as pd
from loader import get_dir
from sending import send
server  = os.getenv('server')
user    = os.getenv('mysqldomain') + '\\' + os.getenv('mysqluser') # Тут правильный двойной слеш!
passwd  = os.getenv('mypassword')
dbase   = os.getenv('db')

eng = sqlalchemy.create_engine(f"mssql+pymssql://{user}:{passwd}@{server}/{dbase}",pool_pre_ping=True)
con = eng.connect()

valume = [
     ['2 группа замечаний: Нет сведений ОМС (Стац.)','Нет сведений ОМС','Стационарная', True]
    ,['2 группа замечаний: Нет сведений ОМС (Амб.)','Нет сведений ОМС','Амбулаторная' , True]
    ,['2 группа замечаний: Нет сведений ОМС (Частные)','Нет сведений ОМС','Частные'   , True]
    #,['4 группа замечаний: Без амбулаторного этапа (Стац.)','Нет амбулаторного этапа','Стационарная']
    ,['4 группа замечаний: Без амбулаторного этапа (Амб.)','Нет амбулаторного этапа','Амбулаторная', False]
    ,['4 группа замечаний: Без амбулаторного этапа (Частные)','Нет амбулаторного этапа','Частные', False]
    ,['5 группа замечаний: Без исхода заболевания больше 30 дней (Стац.)','Без исхода заболевания больше 45 дней','Стационарная', False]
    ,['5 группа замечаний: Без исхода заболевания больше 30 дней (Амб.)','Без исхода заболевания больше 45 дней','Амбулаторная', False]
    ,['5 группа замечаний: Без исхода заболевания больше 30 дней (Частные)','Без исхода заболевания больше 45 дней','Частные', False]
    ,['8 группа замечаний: Неправильный тип лечения (Стац.)', 'Неверный вид лечения', 'Стационарная', False]
    ,['8 группа замечаний: Неправильный тип лечения (Амб.)', 'Неверный вид лечения', 'Амбулаторная',  False]
    ,['8 группа замечаний: Неправильный тип лечения (Частные)', 'Неверный вид лечения', 'Частные', False]
    ,['11 группа замечаний: Количество дублированных  УНРЗ в одном МО (Стац.)', 'Количество дублей','Стационарная', False]
    ,['11 группа замечаний: Количество дублированных  УНРЗ в одном МО (Амб.)', 'Количество дублей','Амбулаторная' , False]
    ,['11 группа замечаний: Количество дублированных  УНРЗ в одном МО (Частные)', 'Количество дублей','Частные'   , False]
    ,['12 группа замечаний: Нет ПАД (Стац.)','Нет ПАД','Стационарная', False]
    ,['12 группа замечаний: Нет ПАД (Амб.)','Нет ПАД','Амбулаторная' , False]
    ,['12 группа замечаний: Нет ПАД (Частные)','Нет ПАД','Частные'   , False]
    ,['13 группа замечаний: Нет дневниковых записей (Стац.)','Нет дневниковых записей','Стационарная', False]
    ,['13 группа замечаний: Нет дневниковых записей (Амб.)','Нет дневниковых записей','Амбулаторная' , False]
    ,['13 группа замечаний: Нет дневниковых записей (Частные)','Нет дневниковых записей','Частные'   , False]
    ,['14 группа замечаний: Зависшие в статусе перевода (Стац.)'  ,'Пациенты зависшие по МО','Стационарная', False]
    ,['14 группа замечаний: Зависшие в статусе перевода (Амб.)'   ,'Пациенты зависшие по МО','Амбулаторная' , False]
    ,['14 группа замечаний: Зависшие в статусе перевода (Частные)','Пациенты зависшие по МО','Частные'   , False]
    ,['15 группа замечаний: Зависшие в статусе перевода без МО прикрепления (Стац.)'  ,'Пациенты зависшие без МО','Стационарная', False]
    ,['15 группа замечаний: Зависшие в статусе перевода без МО прикрепления (Амб.)'   ,'Пациенты зависшие без МО','Амбулаторная' , False]
    ,['15 группа замечаний: Зависшие в статусе перевода без МО прикрепления (Частные)','Пациенты зависшие без МО','Частные'   , False]
]



def generate_pptx(date):
    def create_slide(title,type_error,type_org,date_start,date_end,dynamic):
        if type_org == 'Частные':
            chastn = 'in'
            type_org = 'Амбулаторная'
        else:
            chastn = 'not in'
        
        sql = f"""select distinct case when d1.[Медицинская организация] is null then d2.[Медицинская организация]
                             when d2.[Медицинская организация] is null then d1.[Медицинская организация]
                             else  d1.[Медицинская организация] 
                             end as "Медицинская организация"
                            ,d1.eror1  , d2.eror2 
                            , d2.eror2 -  d1.eror1  as 'Динамика'           
            from (
        SELECT [Медицинская организация]
              , isnull([{type_error}],0) as 'eror1'
          FROM [COVID].[robo].[cv_Zamechania_fr]
          where [дата отчета] = '{date_start}' 
            and [Тип организации] = '{type_org}'  ) as d1
          full join (
        SELECT [Медицинская организация]
              , isnull([{type_error}],0) as 'eror2'
          FROM [COVID].[robo].[cv_Zamechania_fr]
          where [дата отчета] = '{date_end}'
          and [Тип организации] = '{type_org}' ) as d2
          on (d1.[Медицинская организация] = d2.[Медицинская организация])
          where d2.eror2 is not null and d1.eror1 is not null
          and d1.[Медицинская организация] {chastn} ('ООО "ЛАХТА КЛИНИКА"','ООО «Городские поликлиники»', 'ООО "Современная медицина"',
                                                     'ООО «Медицентр ЮЗ»','АНО "Медицинский центр "Двадцать первый век"', 'ООО "ЦСМ "21 ВЕК"',
                                                     'ООО Ава-Петер','ЧУЗ «КБ «РЖД-МЕДИЦИНА» Г. С-ПЕТЕРБУРГ"','ООО «ММЦ «СОГАЗ»')
          and  d1.[Медицинская организация] <> 'МО другого региона'
          order by  [Динамика] DESC, d1.eror1 DESC"""
        
        #send('',sql)

        df = pd.read_sql(sql,con)
        df.fillna(0,inplace=True)
       
        if dynamic:
            df ['Динамика'] = pd.to_numeric(df ['Динамика'])
            df.loc[(df!=0).any(1)]
            df.index = range(len(df))
        else:
            df = df.sort_values(by=['eror2'],ascending=False)
            df = df.loc[df['eror2'] != 0 ]
            df.loc[(df!=0).any(1)]
            df.index = range(len(df)) 
        
        if not len(df):
            return 1
 
        sum_err1 = df['eror1'].sum()
        sum_err2 = df['eror2'].sum()
        title_only_slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(title_only_slide_layout)
        shapes = slide.shapes

        shapes.title.text = title
        shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255,0, 0)
        shapes.title.text_frame.paragraphs[0].font.size = Pt(30)

        body_shape = shapes.placeholders[0]
        tf = body_shape.text_frame
        p = tf.add_paragraph()
        p.text = 'Общее число замечаний на ' + date_start + ' — ' + str(sum_err1).replace('.0','')
        p.font.color.rgb = RGBColor(20,20,20)
        p.font.size = Pt(18)
        
        p = tf.add_paragraph()
        p.text = 'На ' + date_end + ' — ' + str(sum_err2).replace('.0','')
        p.font.color.rgb = RGBColor(20,20,20)
        p.font.size = Pt(18)
        
        p = tf.add_paragraph()
        p.text = 'Общая динамика: ' + str(sum_err2 - sum_err1).replace('.0','')
        p.font.color.rgb = RGBColor(20,20,20)
        p.font.size = Pt(18)
        
        # Первая таблица  
        if dynamic:
            if len(df.loc[~(df['Динамика'] < 0)]) > 16:
                rows = 17
            else:
                rows = len(df.loc[~(df['Динамика'] < 0)]) + 1
        else:
            #send('', type_error + str(len(df)))
            if len(df) > 16:
                rows = 17
            else:
                rows = len(df) + 1


        if dynamic:
            cols = 5
        else:
            cols = 3
        left = Inches(0.5)
        top = Inches(2)
        width = Inches(5)
        height = Inches(0.1)

        shape = shapes.add_table(rows, cols, left, top, width, height)
        table = shape.table
        tbl = shape._element.graphic.graphicData.tbl
        style_id = '{9DCAF9ED-07DC-4A11-8D7F-57B35C25682E}'
        tbl[0][-1].text = style_id

        # set column widths
        table.columns[0].width = Inches(0.5)
        table.columns[1].width = Inches(3.5)
        table.columns[2].width = Inches(1.5)
        if dynamic:
            table.columns[1].width = Inches(2.5)
            table.columns[3].width = Inches(1.5)
            table.columns[4].width = Inches(1.5)
        # write column headings
        table.cell(0, 0).text = '№'
        table.cell(0, 1).text = 'Медицинская организация'
        if dynamic:
            table.cell(0, 2).text = date_start
            table.cell(0, 3).text = date_end
            table.cell(0, 4).text = 'Динамика'
        else:
            table.cell(0, 2).text = date_end

        k = 0
        if dynamic:
            if len(df.loc[~(df['Динамика'] < 0)]):
                for i in df.loc[~(df['Динамика'] < 0)].head(16).index:
                    table.cell(k+1, 0).text = str(i+1) 
                    table.cell(k+1, 0).text_frame.paragraphs[0].font.size = Pt(12)
                    table.cell(k+1, 1).text = df.at[i,'Медицинская организация']
                    table.cell(k+1, 1).text_frame.paragraphs[0].font.size = Pt(10) 
                    table.cell(k+1, 2).text = str(df.at[i,'eror1']).replace('.0','')
                    table.cell(k+1, 2).text_frame.paragraphs[0].font.size = Pt(12)
                    table.cell(k+1, 3).text = str(df.at[i,'eror2']).replace('.0','')
                    table.cell(k+1, 3).text_frame.paragraphs[0].font.size = Pt(12)
                    table.cell(k+1, 4).text = str(df.at[i,'Динамика']).replace('.0','')
                    table.cell(k+1, 4).text_frame.paragraphs[0].font.size = Pt(12)
                    k += 1
                    if k > 16:
                        break
        else:
            for i in df.head(16).index:
                table.cell(k+1, 0).text = str(k+1)
                table.cell(k+1, 0).text_frame.paragraphs[0].font.size = Pt(12)
                table.cell(k+1, 1).text = df.at[i,'Медицинская организация']
                table.cell(k+1, 1).text_frame.paragraphs[0].font.size = Pt(10)
                table.cell(k+1, 2).text = str(df.at[i,'eror2']).replace('.0','')
                table.cell(k+1, 2).text_frame.paragraphs[0].font.size = Pt(12)
                df.drop(i,inplace=True)
                k += 1

        # Вторая таблица  
        if dynamic:
            if len(df.loc[df['Динамика'] < 0]) > 16:
                rows = 17
            else:
                rows = len(df.loc[df['Динамика'] < 0]) + 1
        else:
            #send('','2' +  type_error + str(len(df)))
            if len(df) > 16:
                rows = 17
            else:
                rows = len(df) + 1
                if rows < 2:
                    #send('', str(df.head()))
                    return 1


        if dynamic:
            cols = 5
        else:
            cols = 3
        left = Inches(8.25)
        top = Inches(2)
        width = Inches(5)
        height = Inches(0.1)

        shape = shapes.add_table(rows, cols, left, top, width, height)
        table = shape.table
        tbl = shape._element.graphic.graphicData.tbl
        style_id = '{9DCAF9ED-07DC-4A11-8D7F-57B35C25682E}'
        tbl[0][-1].text = style_id

        # set column widths
        table.columns[0].width = Inches(0.5)
        table.columns[1].width = Inches(3.5)
        table.columns[2].width = Inches(1.5)
        if dynamic:
            table.columns[1].width = Inches(2.5)
            table.columns[3].width = Inches(1.5)
            table.columns[4].width = Inches(1.5)
        # write column headings
        table.cell(0, 0).text = '№'
        table.cell(0, 1).text = 'Медицинская организация'
        if dynamic:
            table.cell(0, 2).text = date_start 
            table.cell(0, 3).text = date_end
            table.cell(0, 4).text = 'Динамика'
        else:
            table.cell(0, 2).text = date_end

        k = 0
        if dynamic:
            if len(df.loc[df['Динамика'] < 0]):
                for i in df.loc[df['Динамика'] < 0].sort_values(by=['Динамика']).head(16).index:
                    table.cell(k+1, 0).text = str(i+1) 
                    table.cell(k+1, 0).text_frame.paragraphs[0].font.size = Pt(12)
                    table.cell(k+1, 1).text = df.at[i,'Медицинская организация']
                    table.cell(k+1, 1).text_frame.paragraphs[0].font.size = Pt(10) 
                    table.cell(k+1, 2).text = str(df.at[i,'eror1']).replace('.0','')
                    table.cell(k+1, 2).text_frame.paragraphs[0].font.size = Pt(12)
                    table.cell(k+1, 3).text = str(df.at[i,'eror2']).replace('.0','')
                    table.cell(k+1, 3).text_frame.paragraphs[0].font.size = Pt(12)
                    table.cell(k+1, 4).text = str(df.at[i,'Динамика']).replace('.0','')
                    table.cell(k+1, 4).text_frame.paragraphs[0].font.size = Pt(12)
                    k += 1
                    if k > 16:
                        break
        else:
            for i in df.head(16).index:
                table.cell(k+1, 0).text = str(k+17)
                table.cell(k+1, 0).text_frame.paragraphs[0].font.size = Pt(12)
                table.cell(k+1, 1).text = df.at[i,'Медицинская организация']
                table.cell(k+1, 1).text_frame.paragraphs[0].font.size = Pt(10)
                table.cell(k+1, 2).text = str(df.at[i,'eror2']).replace('.0','')
                table.cell(k+1, 2).text_frame.paragraphs[0].font.size = Pt(12)
                k += 1

    date1 = date.split(',')[0]
    date2 = date.split(',')[1]
    date_start_sql = date1
    date_end_sql = date2
    date_start = pd.to_datetime(date1).strftime("%d.%m.%Y")
    date_end = pd.to_datetime(date2).strftime("%d.%m.%Y")
    prs = Presentation()
    
    prs.slide_height = Inches(9)
    prs.slide_width = Inches(16)

    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    # Первый слайд
    title.text = "Замечания по ведению Федерального регистра лиц, больных COVID-19"
    subtitle.text = "По состоянию на " + date_end

    # Второй слайд


    bullet_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[0]

    title_shape.text = 'Группы замечаний по ведению Регистра на '+ date_end
    tf = body_shape.text_frame
    date_in_text = (pd.to_datetime(date2) - datetime.timedelta(30)).strftime("%d.%m.%Y")

    text = f"""Срок создания регистровой записи (УНРЗ) не соответствует дате установки диагноза – более 7 дней между датами (МЗ оценивает регион) (Количество регистровых записей, внесённых с задержкой больше недели, за последний месяц)
    Отсутствие информации о номере Полиса ОМС в разделе «Медицинское страхование»
    Не заполнен Раздел «Результаты ежедневного наблюдения» (дневниковые записи) – по данным МЗ РФ
    Поле «исход заболевания» заполнено «переведён в другую МО» без открытия следующего этапа лечения («зависшие» пациенты более 7 дней от даты перевода в др. МО)
    Поле «исход заболевания» не заполнено с датой установки диагноза ранее {date_in_text} (длительность болезни более 30 дней)
    Количество УНРЗ с заполненным полем «исход заболевания» «выздоровление» не соответствует оперативной отчётности поликлиник
    Количество УНРЗ с пустым полем «исход заболевания» (находящиеся под медицинским наблюдением в амбулаторных условиях) не соответствует оперативной отчётности поликлиник
    Поле «Вид лечения» ошибочно заполнено «стационарный» -  у пациентов, находящихся под медицинским наблюдением в амбулаторных условиях 
    Количество УНРЗ с заполненным Разделом «Результаты ежедневного наблюдения» в части поля «ИВЛ» не соответствует оперативной отчётности стационаров
    Количество УНРЗ пациентов, находящихся под медицинским наблюдением в стационарных условиях, не соответствует оперативной отчётности стационаров
    Дублированные  УНРЗ в одном МО
    Отсутствие прикреплённого файла патологоанатомического заключения"""

    rows = 12
    cols = 2 
    left = Inches(0.5)
    top = Inches(2)
    width = Inches(8)
    height = Inches(0.1)

    shape = shapes.add_table(rows, cols, left, top, width, height)

    tbl = shape._element.graphic.graphicData.tbl
    table = shape.table

    style_id = '{2D5ABB26-0587-4C30-8999-92F81FD0307C}'
    tbl[0][-1].text = style_id

    table.columns[0].width = Inches(0.5)
    table.columns[1].width = Inches(8.0)

    i = 0
    for t in text.split('\n'):
        table.cell(i, 0).text = str(i+1) + ')' 
        table.cell(i, 0).text_frame.paragraphs[0].font.size = Pt(12)
        table.cell(i, 1).text = t
        table.cell(i, 1).text_frame.paragraphs[0].font.size = Pt(11)
        i+=1
        
    # Третий слайд

    bullet_slide_layout = prs.slide_layouts[2]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes


    title_shape = shapes.title
    title_shape.text ='1 группа замечаний'
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255,0, 0)

    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame


    #df = pd.read_sql("Select count(*) from cv_fedreg where DATEDIFF(day,[Дата создания РЗ],getdate() ) < 31", con)
    df = pd.read_sql("""
	    select count(*)
	    from cv_fedreg 
	    where 
	    DATEDIFF(day,[Дата создания РЗ],getdate() ) < 31
	    and (DATEDIFF(day,[Дата создания РЗ],[Диагноз установлен] ) > 7 or DATEDIFF(day,[Дата создания РЗ],[Диагноз установлен] ) < -7 )
            """, con)


    p = tf.add_paragraph()
    p.text = 'за последние 30 дней: + ' + str(df.iat[0,0]) + ' УНРЗ'
    p.font.color.rgb = RGBColor(255,0,0)
    p.font.size = Pt(30)

    df = pd.read_sql("""
	    select count(*)
	    from cv_fedreg 
	    where 
	    (DATEDIFF(day,[Дата создания РЗ],[Диагноз установлен] ) > 7 or DATEDIFF(day,[Дата создания РЗ],[Диагноз установлен] ) < -7 )
            """, con)

    p = tf.add_paragraph()
    p.text = 'за всё время: ' + str(df.iat[0,0]) + ' УНРЗ'
    p.font.color.rgb = RGBColor(255,0,0)
    p.font.size = Pt(30)

    p = tf.add_paragraph()
    p = tf.add_paragraph()
    p.text = 'Срок создания регистровой записи (УНРЗ) не соответствует дате установки диагноза: разница более 7 дней после даты установки диагноза, количество взято за последний месяц'
    p.font.size = Pt(20)

    for title,type_error,type_org,dynamic in valume:
        try:
            create_slide(title,type_error,type_org,date_start_sql,date_end_sql,dynamic)
        except Exception as e:
            send('', type_error +'\n'+ type_org +'\n'+ str(e))

    pptx_file = get_dir('temp') + '/dynamic.pptx' 
    prs.save(pptx_file)

    return pptx_file
