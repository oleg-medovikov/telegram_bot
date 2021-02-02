import os,datetime,pyodbc
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Inches
import pandas as pd

conn = pyodbc.connect(os.getenv('sql_conn'))

valume = [
     ['2 группа замечаний: Нет сведений ОМС (Стац.)','Нет сведений ОМС','Стационарная']
    ,['2 группа замечаний: Нет сведений ОМС (Амб.)','Нет сведений ОМС','Амбулаторная']
    #,['4 группа замечаний: Без амбулаторного этапа (Стац.)','Нет амбулаторного этапа','Стационарная']
    ,['4 группа замечаний: Без амбулаторного этапа (Амб.)','Нет амбулаторного этапа','Амбулаторная']
    ,['5 группа замечаний: Без исхода заболевания больше 45 дней (Стац.)','Без исхода заболевания больше 45 дней','Стационарная']
    ,['5 группа замечаний: Без исхода заболевания больше 45 дней (Амб.)','Без исхода заболевания больше 45 дней','Амбулаторная']
    ,['8 группа замечаний: Неправильный тип лечения (Стац.)', 'Неверный вид лечения', 'Стационарная']
    ,['8 группа замечаний: Неправильный тип лечения (Амб.)', 'Неверный вид лечения', 'Амбулаторная']
    ,['11 группа замечаний: Количество дублированных  УНРЗ в одном МО (Стац.)', 'Количество дублей','Стационарная']
    ,['11 группа замечаний: Количество дублированных  УНРЗ в одном МО (Амб.)', 'Количество дублей','Амбулаторная']
    ,['12 группа замечаний: Нет ПАД (Стац.)','Нет ПАД','Стационарная' ]
    ,['12 группа замечаний: Нет ПАД (Амб.)','Нет ПАД','Амбулаторная' ]
    
]



def generate_pptx(date):
    def create_slide(title,type_error,type_org,date_start):
        sql = f"""
        select distinct d1.[Медицинская организация] ,d1.eror1  , d2.eror2 
                    , d1.eror1 -  d2.eror2 as 'Динамика'           
            from (
        SELECT [Медицинская организация]
              , [{type_error}] as 'eror1'
          FROM [COVID].[robo].[cv_Zamechania_fr]
          where [дата отчета] = '{date_start}' 
            and [Тип организации] = '{type_org}'  ) as d1
          full join (
        SELECT [Медицинская организация]
              , [{type_error}] as 'eror2'
          FROM [COVID].[robo].[cv_Zamechania_fr]
          where [дата отчета] = cast(getdate() as date)
          and [Тип организации] = '{type_org}' ) as d2
          on (d1.[Медицинская организация] = d2.[Медицинская организация])
          order by  d1.eror1 - d2.eror2  DESC , d1.eror1 DESC
        """
        df = pd.read_sql(sql,conn)
        
        sum_err1 = df['eror1'].sum()
        sum_err2 = df['eror2'].sum()
        
        title_only_slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(title_only_slide_layout)
        shapes = slide.shapes

        shapes.title.text = title
        shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255,0, 0)
        shapes.title.text_frame.paragraphs[0].font.size = Pt(30)

        body_shape = shapes.placeholders[0]
        tf = body_shape.text_frame
        p = tf.add_paragraph()
        p.text = 'Общее число замечаний на дату ' + date_start + ' было ' + str(sum_err1)
        p.font.color.rgb = RGBColor(20,20,20)
        p.font.size = Pt(18)
        
        p = tf.add_paragraph()
        p.text = 'А на дату ' + date_end + ' cтало ' + str(sum_err2)
        p.font.color.rgb = RGBColor(20,20,20)
        p.font.size = Pt(18)
        
        p = tf.add_paragraph()
        p.text = 'Общая динамика: ' + str(sum_err2 - sum_err1)
        p.font.color.rgb = RGBColor(20,20,20)
        p.font.size = Pt(18)
        
        rows = 17
        cols = 5
        left = Inches(0.5)
        top = Inches(2)
        width = Inches(5)
        height = Inches(0.1)

        table = shapes.add_table(rows, cols, left, top, width, height).table

        # set column widths
        table.columns[0].width = Inches(0.5)
        table.columns[1].width = Inches(4.0)
        table.columns[2].width = Inches(1.5)
        table.columns[3].width = Inches(1.5)
        table.columns[4].width = Inches(1.5)
        # write column headings
        table.cell(0, 0).text = '№'
        table.cell(0, 1).text = 'Медицинская организация'
        table.cell(0, 2).text = date_start
        table.cell(0, 3).text = date_end
        table.cell(0, 4).text = 'Динамика'
        for i in range(16):
            table.cell(i+1, 0).text = str(i+1) 
            table.cell(i+1, 0).text_frame.paragraphs[0].font.size = Pt(12)
            table.cell(i+1, 1).text = df.at[i,'Медицинская организация']
            table.cell(i+1, 1).text_frame.paragraphs[0].font.size = Pt(10) 
            table.cell(i+1, 2).text = str(df.at[i,'eror1'])
            table.cell(i+1, 2).text_frame.paragraphs[0].font.size = Pt(12)
            table.cell(i+1, 3).text = str(df.at[i,'eror2'])
            table.cell(i+1, 3).text_frame.paragraphs[0].font.size = Pt(12)
            table.cell(i+1, 4).text = str(df.at[i,'Динамика'])
            table.cell(i+1, 4).text_frame.paragraphs[0].font.size = Pt(12)

    date_start = date
    date_end = datetime.datetime.today().strftime("%Y-%m-%d")
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    # Первый слайд
    title.text = "Замечания по ведению Федерального регистра лиц, больных COVID-19"
    subtitle.text = "По состоянию на " + datetime.datetime.today().strftime("%Y-%m-%d")

    try:
        os.remove(os.getenv('path_temp')+r'\test.pptx')
    except:
        pass
    # Второй слайд


    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Группы замечаний по ведению Регистра на '+ datetime.datetime.today().strftime("%Y-%m-%d")
    tf = body_shape.text_frame

    def paragraf(text,r,g,b):
        p = tf.add_paragraph()
        p.text = text
        p.font.color.rgb = RGBColor(r,g,b)
        p.font.size = Pt(12)

    text = '1) Срок создания регистровой записи (УНРЗ) не соответствует дате установки диагноза – более 7 дней между датами (МЗ оценивает регион) (Количество регистровых записей, внесенных с задержкой больше недели, за последний месяц)'
    paragraf(text,255,0,0)
    text = """2) Отсутствие информации о номере Полиса ОМС в разделе «Медицинское страхование»
    3) Не заполнен Раздел «Результаты ежедневного наблюдения» (дневниковые записи) – по данным МЗ РФ
    4) Поле «исход заболевания» заполнено «переведен в другую МО» без открытия следующего этапа лечения («зависшие» пациенты более 7 дней от даты перевода в др. МО)
    5) Поле «исход заболевания» не заполнено с датой установки диагноза ранее 15.12.2020  (длительность болезни более 45 дней)
    6) Количество УНРЗ с заполненным полем «исход заболевания» «выздоровление» не соответствует оперативной отчетности поликлиник
    7) Количество УНРЗ с пустым полем «исход заболевания» (находящиеся под медицинским наблюдением в амбулаторных условиях) не соответствует оперативной отчетности поликлиник
    8) Поле «Вид лечения» ошибочно заполнено «стационарный» -  у пациентов, находящихся под медицинским наблюдением в амбулаторных условиях 
    9) Количество УНРЗ с заполненным Разделом «Результаты ежедневного наблюдения» в части поля «ИВЛ» не соответствует оперативной отчетности стационаров
    10) Количество УНРЗ пациентов, находящихся под медицинским наблюдением в стационарных условиях, не соответствует оперативной отчетности стационаров
    11) Дублированные  УНРЗ в одном МО
    12) Отсутствие прикрепленного файла патологоанатомического заключения
    """

    for t in text.split('\n'):
        paragraf(t,0,0,0)
        
    # Третий слайд

    bullet_slide_layout = prs.slide_layouts[2]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes


    title_shape = shapes.title
    title_shape.text ='1 группа замечаний'
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255,0, 0)

    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame


    df = pd.read_sql("Select count(*) from cv_fedreg where DATEDIFF(day,[Дата создания РЗ],getdate() ) < 31", conn)

    p = tf.add_paragraph()
    p.text = 'за последние 30 дней: + ' + str(df.iat[0,0]) + ' УНРЗ'
    p.font.color.rgb = RGBColor(255,0,0)
    p.font.size = Pt(30)

    p = tf.add_paragraph()
    p = tf.add_paragraph()
    p = tf.add_paragraph()
    p.text = 'Срок создания регистровой записи (УНРЗ) не соответствует дате установки диагноза: разница более 7 дней после даты установки диагноза, количество взято за последний месяц'
    p.font.size = Pt(20)

    for title,type_error,type_org in valume:
        create_slide(title,type_error,type_org,date_start)

    prs.save(os.getenv('path_temp')+r'\test.pptx')
    
    return(os.getenv('path_temp')+r'\test.pptx')
